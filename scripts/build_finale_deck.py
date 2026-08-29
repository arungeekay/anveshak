"""Build the Grand Finale pitch deck (FINALE_PLAN F-18).

Different job from the submission deck: that one filled a form, this one supports a
six-minute *performance*. So it is short, the live demo is the centrepiece, and
every slide exists to set up or land one claim.

Uses the same 16:9 template as the submission deck, and the same hard-won lesson:
the template's content area renders WHITE, so text must be dark (see build_deck.py).

    python scripts/build_finale_deck.py
Out: ANVESHAK_Finale_Pitch.pptx (+ .pdf via PowerPoint, if available)
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

TEMPLATE = "KSP Datathon 2026 _ Prototype Submission Template (1).pptx"
OUT = "ANVESHAK_Finale_Pitch.pptx"
LIVE = "https://anveshak-api-50044329134.development.catalystappsail.in/ui/"

INK = RGBColor(0x0B, 0x1F, 0x3A)      # headings / emphasis
BODY = RGBColor(0x2B, 0x37, 0x4B)     # body copy
ACCENT = RGBColor(0x1D, 0x4E, 0xD8)   # metrics / links
MUTED = RGBColor(0x64, 0x74, 0x8B)    # captions
GOOD = RGBColor(0x04, 0x78, 0x55)     # confirmed / live
WARN = RGBColor(0xB4, 0x53, 0x09)     # the cost of not having it

BODY_BOX = (0.55, 1.55, 8.95, 3.75)   # left, top, width, height (inches)
EMU = 914400


def add_body(slide):
    """Trim the template's prompt to its title and return a fresh body frame."""
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            for p in list(sh.text_frame.paragraphs)[1:]:
                p._p.getparent().remove(p._p)
            break
    left, top, width, height = BODY_BOX
    tb = slide.shapes.add_textbox(Emu(int(left * EMU)), Emu(int(top * EMU)),
                                  Emu(int(width * EMU)), Emu(int(height * EMU)))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def set_title(slide, text: str) -> None:
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            p = sh.text_frame.paragraphs[0]
            for r in list(p.runs):
                r._r.getparent().remove(r._r)
            r = p.add_run()
            r.text = text
            r.font.bold = True
            r.font.size = Pt(24)
            r.font.color.rgb = INK
            return


def para(tf, text, *, size=13, color=BODY, bold=False, space=7, first=False,
         bullet=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if first:
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
    p.space_after = Pt(space)
    runs = text if isinstance(text, list) else [(text, color, bold)]
    for i, item in enumerate(runs):
        t, c, b = item if isinstance(item, tuple) else (item, color, bold)
        r = p.add_run()
        r.text = ("• " if bullet and i == 0 else "") + t
        r.font.size = Pt(size)
        r.font.color.rgb = c
        r.font.bold = b
        r.font.name = "Calibri"
    return p


def keep(prs, indices):
    """Keep only these template slides, in order."""
    lst = prs.slides._sldIdLst
    ids = list(lst)
    for i, sid in enumerate(ids):
        if i not in indices:
            lst.remove(sid)


def main() -> int:
    prs = Presentation(TEMPLATE)
    # Reuse 8 content slides plus the official THANK YOU closer. keep() renumbers
    # what survives, so index S by POSITION AFTER the cull, not by template index.
    last = len(prs.slides._sldIdLst) - 1
    keep(prs, {0, 1, 2, 3, 6, 11, 12, 13, last})
    S = list(prs.slides)   # 0 title · 1 problem · 2 what · 3 demo · 4 counterfactual
                           # 5 trust · 6 architecture · 7 roadmap · 8 thank you

    # 1 — Title / team
    box = next(s for s in S[0].shapes if s.has_text_frame)
    tf = box.text_frame
    tf.clear()
    para(tf, "ANVESHAK · ಅನ್ವೇಷಕ", size=22, color=INK, bold=True, first=True)
    para(tf, "An AI that doesn't just answer questions — it works cases.",
         size=15, color=BODY)
    para(tf, "Team Zen — Hiran Vikraman S R · Arun G K   |   Challenge 1",
         size=12, color=MUTED)
    para(tf, f"Live: {LIVE}", size=11, color=ACCENT)

    # 2 — The officer's problem
    set_title(S[1], "One night in Jayanagar")
    tf = add_body(S[1])
    para(tf, "A woman's gold chain is snatched by two men on a motorcycle. "
             "The FIR is filed at Jayanagar station.", size=14, color=INK,
         bold=True, first=True, space=10)
    para(tf, "The same men struck in Mandya last month, and Tumakuru before that.",
         bullet=True)
    para(tf, "Nobody knows. The records sit in three district silos, and no one "
             "has time to read 15,000 narratives looking for a pattern.", bullet=True)
    para(tf, "To even ask 'how many chain snatchings this year?', an officer needs "
             "someone who writes SQL — and English.", bullet=True)
    para(tf, [("This is the gap ANVESHAK closes.", INK, True)], size=14, space=2)

    # 3 — What it does
    set_title(S[2], "Ask. Then let it work the case.")
    tf = add_body(S[2])
    for h, rest in [
        ("Ask in Kannada or English", " — voice or text, and every answer arrives "
         "with its SQL, its rows and its case IDs."),
        ("Find the ring", " — MO fingerprinting links crimes across district "
         "boundaries that siloed records never connect."),
        ("Work the case", " — six specialist agents assemble a court-ready "
         "Investigation Pack in about ten seconds."),
        ("Patrol tonight", " — an overnight sweep raises ranked leads and turns "
         "them into a deployment plan."),
        ("Prove it", " — role-scoped access, a tamper-evident audit log, and "
         "guardrails you are welcome to attack."),
    ]:
        para(tf, [(h, INK, True), (rest, BODY, False)], size=13, bullet=True,
             first=(h.startswith("Ask")))

    # 4 — LIVE DEMO
    set_title(S[3], "Live demo")
    tf = add_body(S[3])
    para(tf, "From the FIR a constable files, to the pack an IO takes to court.",
         size=16, color=INK, bold=True, first=True, space=12)
    for step in [
        "The overnight sweep — and the morning briefing, in Kannada",
        "A question, answered with its evidence",
        "Try to break it: injection, exfiltration, profiling",
        "File a new FIR — in your own words — and watch it join the series",
        "Replay the ring crossing three districts",
        "Six agents work the case → the pack → the suspect",
        "Tonight's patrol plan · the same question as SHO, and as Analyst",
    ]:
        para(tf, step, size=12.5, bullet=True, space=5)

    # 5 — The cost of not having it
    set_title(S[4], "What it would have changed")
    tf = add_body(S[4])
    para(tf, "SH-07 · Operation Gold Chain Black Visor", size=15, color=INK,
         bold=True, first=True, space=10)
    para(tf, [("ANVESHAK flags this ring at case #6, on 10 February 2026.",
               INK, True)], size=14, space=8)
    para(tf, [("Nine further chain snatchings", WARN, True),
              (" across Bengaluru City, Mandya and Tumakuru followed over the next ",
               BODY, False), ("142 days", WARN, True),
              (" — after the pattern was already visible.", BODY, False)],
         size=14, space=10)
    para(tf, "Computed by replaying our own linkage engine over the corpus "
             "truncated at each of the series' case dates. A retrospective replay "
             "on synthetic data — not a prediction, and not a claim about a real "
             "case.", size=11, color=MUTED, space=2)

    # 6 — Why you can trust it
    set_title(S[5], "Built to be checked, not believed")
    tf = add_body(S[5])
    for h, rest in [
        ("The model never computes a number", " — deterministic tools do; the model "
         "only narrates. Every figure traces to a tool result and a case ID."),
        ("Attack it yourself", " — a red-team console in the product. 10 of 10 "
         "attack vectors blocked, re-tested on every page load."),
        ("Religion and caste are never model features", " — a profiling question is "
         "refused, in English and Kannada, and the refusal is logged."),
        ("Scope is enforced in the query", " — SHO sees a station, SP a district, "
         "Analyst sees masked names. Not hidden in the browser."),
        ("History cannot be rewritten", " — the audit log is hash-chained; press "
         "'Verify chain' and watch it recompute."),
        ("Our ground truth is public", " — linkage precision 0.86, recall 12/14, "
         "checkable against the planted data in the repository."),
    ]:
        para(tf, [(h, INK, True), (rest, BODY, False)], size=12, bullet=True,
             space=5, first=h.startswith("The model"))

    # 7 — Architecture
    set_title(S[6], "Runs entirely on Zoho Catalyst")
    tf = add_body(S[6])
    para(tf, "Officer (EN/KN, voice or text) → React SPA → FastAPI on AppSail",
         size=13, color=INK, bold=True, first=True, space=8)
    for line in [
        "intent router → NL→SQL (guardrails · ADR-9 policy · role scope · self-repair)",
        "            → typed tools: linkage · graph · forecast · hotspots · risk · similar",
        "DuckDB analytical mirror ⟷ Catalyst Data Store (system of record)",
        "Catalyst QuickML → GLM-4.7-Flash   ·   SmartBrowz → pack PDF",
        "Cron → Night Patrol → lead cards → morning briefing",
    ]:
        para(tf, line, size=12, color=BODY, space=4)
    para(tf, "No external inference API is ever called from the deployed app.",
         size=12, color=GOOD, bold=True, space=2)

    # 8 — Roadmap / ask
    set_title(S[7], "What we would build next")
    tf = add_body(S[7])
    for line in [
        "CCTNS integration: live read connection + identity resolution on "
        "name / parentage / DOB, replacing our synthetic person key.",
        "Catalyst Auth issuing real officer identities into the scope layer that "
        "already enforces them.",
        "Kannada NL→SQL tuning with hand-curated examples from real officer "
        "phrasing — our clearest measured gap.",
        "Human-in-the-loop learning: every Confirm/Reject on a series becomes a "
        "label that improves precision.",
        "A pilot with one district's live data — that is the ask.",
    ]:
        para(tf, line, size=12.5, bullet=True, space=6,
             first=line.startswith("CCTNS"))

    prs.save(OUT)
    print(f"wrote {OUT}  ({len(prs.slides._sldIdLst)} slides)")
    print("PDF: open in PowerPoint → Save as PDF, or run the export in build_deck.py")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
