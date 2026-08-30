"""Fill the official KSP Datathon 2026 submission template with ANVESHAK content.

The template has a full-slide BLACK background image + a bold heading text box per
slide (the section prompt). We keep each heading as the slide title and add a light
(readable-on-black) body text box beneath it. All numbers are the measured ones from
README.md / eval / PROGRESS.log. No invented figures.

Run:  python scripts/build_deck.py
Out:  ANVESHAK_KSP_Datathon_2026_Submission.pptx
"""
from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

TEMPLATE = "KSP Datathon 2026 _ Prototype Submission Template (1).pptx"
OUT = "ANVESHAK_KSP_Datathon_2026_Submission.pptx"

LIVE_URL = "https://anveshak-api-50044329134.development.catalystappsail.in/ui/"
API_URL = "https://anveshak-api-50044329134.development.catalystappsail.in"
GITHUB = "https://github.com/arungeekay/anveshak"

# Palette tuned for the template's WHITE content background (dark text on white).
# (Names kept for minimal churn; "WHITE" now means the dark emphasis ink.)
WHITE = RGBColor(0x0B, 0x1F, 0x3A)    # dark navy: headings / emphasis
BODY = RGBColor(0x2B, 0x37, 0x4B)     # dark slate: body copy
ACCENT = RGBColor(0x1D, 0x4E, 0xD8)   # blue-700: metrics / links
MUTED = RGBColor(0x64, 0x74, 0x8B)    # slate-500: captions / placeholders
GOOD = RGBColor(0x04, 0x78, 0x55)     # emerald-700: live / confirmed

# body text box geometry (below the heading), in inches
BODY_LEFT, BODY_TOP, BODY_W, BODY_H = 0.55, 1.55, 8.95, 3.75


def add_body(slide):
    # The template's prompt box sometimes carries multi-line guidance (e.g. the
    # Opportunities sub-questions, the Links numbered list) that would sit right where
    # our body goes. Trim that box to just its first line (the section title) so our
    # body never collides with leftover prompt text.
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            for p in list(sh.text_frame.paragraphs)[1:]:
                p._p.getparent().remove(p._p)
            break
    tb = slide.shapes.add_textbox(
        Emu(int(BODY_LEFT * 914400)), Emu(int(BODY_TOP * 914400)),
        Emu(int(BODY_W * 914400)), Emu(int(BODY_H * 914400)))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def para(tf, text, *, size=12, color=BODY, bold=False, level=0, bullet=False,
         space=4, first=False):
    if first:
        p = tf.paragraphs[0]
        for r in list(p.runs):  # clear any existing runs on the reused paragraph
            r._r.getparent().remove(r._r)
    else:
        p = tf.add_paragraph()
    p.level = level
    p.space_after = Pt(space)
    runs = text if isinstance(text, list) else [(text, color, bold)]
    for i, item in enumerate(runs):
        t, c, b = item if isinstance(item, tuple) else (item, color, bold)
        r = p.add_run()
        r.text = ("• " if bullet and i == 0 else "") + t
        r.font.size = Pt(size)
        r.font.color.rgb = c
        r.font.bold = b
        r.font.name = "Calibri"
    return p


SHOTS = [
    ("video/shots/01_leadfeed.png", "Night Patrol: proactive leads"),
    ("video/shots/02_chat_evidence.png", "Chat: answer plus its SQL evidence"),
    ("video/shots/07_intake.png", "New FIR joins a live series"),
    ("video/shots/04_pack.png", "Investigation Pack: six agents"),
    ("video/shots/08_trust.png", "Trust Center: red-team console"),
    ("video/shots/09_person.png", "Person 360"),
]


def _place_shots(slide):
    """Lay the six live-app screenshots in a 3x2 grid with captions under each."""
    import os as _os
    EMU = 914400
    cols, gap = 3, 0.18
    left0, top0 = 0.35, 1.48
    iw = (9.3 - gap * (cols - 1)) / cols          # image width in inches
    ih = iw * 9 / 16                               # 16:9
    cap_h, row_gap = 0.26, 0.16
    for i, (path, caption) in enumerate(SHOTS):
        if not _os.path.exists(path):
            continue
        r, c = divmod(i, cols)
        x = left0 + c * (iw + gap)
        y = top0 + r * (ih + cap_h + row_gap)
        slide.shapes.add_picture(path, Emu(int(x * EMU)), Emu(int(y * EMU)),
                                 Emu(int(iw * EMU)), Emu(int(ih * EMU)))
        tb = slide.shapes.add_textbox(Emu(int(x * EMU)), Emu(int((y + ih + 0.01) * EMU)),
                                      Emu(int(iw * EMU)), Emu(int(cap_h * EMU)))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r0 = p.add_run()
        r0.text = caption
        r0.font.size = Pt(9.5)
        r0.font.color.rgb = BODY
        r0.font.name = "Calibri"


def _delete_slide_by_text(prs, marker):
    """Remove the first slide whose text equals `marker` (e.g. the 'Blank slide')."""
    for i, s in enumerate(prs.slides):
        if any(sh.has_text_frame and sh.text_frame.text.strip() == marker for sh in s.shapes):
            lst = prs.slides._sldIdLst
            lst.remove(list(lst)[i])
            return True
    return False


def fill_team(slide):
    """Slide 0: replace the Team Details prompt lines with our values."""
    box = next(s for s in slide.shapes if s.has_text_frame and "Team Details" in s.text_frame.text)
    tf = box.text_frame
    lines = [
        ("Team Details", WHITE, True, 17),
        ("", BODY, False, 6),
        ("Team name:  Team Zen", BODY, False, 14),
        ("Team leader name:  Hiran Vikraman S R", BODY, False, 14),
        ("Team size:  2   (Hiran Vikraman S R · Arun G K)", BODY, False, 14),
        ("Problem Statement:  Challenge 1. Intelligent Conversational AI for the "
         "KSP Crime Database. Enable officers to query the FIR/crime database in "
         "natural language (English and Kannada) and get accurate, evidence-grounded "
         "answers.", BODY, False, 13),
    ]
    tf.clear()  # leaves a single empty paragraph
    for i, (text, color, bold, size) in enumerate(lines):
        para(tf, text, size=size, color=color, bold=bold, first=(i == 0))


def main():
    prs = Presentation(TEMPLATE)
    S = prs.slides

    fill_team(S[0])

    # Slide 1: Brief about the solution
    tf = add_body(S[1])
    para(tf, [("ANVESHAK (ಅನ್ವೇಷಕ", WHITE, True),
              (' = "the investigator") turns the KSP crime database into an ', BODY, False),
              ("AI investigation partner", WHITE, True),
              (", not just a chatbot.", BODY, False)], size=13, first=True, space=8)
    para(tf, [("Ask in Kannada or English, by voice or text, and get an answer in "
               "seconds ", BODY, False), ("with the evidence attached", WHITE, True),
              (": the SQL, the rows, the case IDs. No SQL skills, no analyst queue.",
               BODY, False)], bullet=True)
    para(tf, [("Deterministic tools compute; the model only narrates.", WHITE, True),
              (" Every number traces to a tool result, so there are no invented "
               "statistics. That is what makes it usable in policing.", BODY, False)],
         bullet=True)
    para(tf, "It does not stop at answers, it works cases: links serial crime across "
             "district boundaries, files new FIRs and links them live, runs a six-agent "
             "cell that assembles court-ready packs, and sweeps overnight to raise leads "
             "before anyone asks.", bullet=True)
    para(tf, [("Impact: minutes instead of days, catches that siloed records miss, and "
               "a ranked lead list every morning. ", WHITE, True),
              ("Runs 100% on Zoho Catalyst.", GOOD, True)], bullet=True, space=2)

    # Slide 2: Opportunities / differentiation / USP
    tf = add_body(S[2])
    para(tf, "How it is different, and how it solves the problem:", size=12, color=WHITE,
         bold=True, first=True, space=8)
    para(tf, 'Most "crime chatbots" stop at natural language to SQL. ANVESHAK adds an '
             "investigation layer: cross-district serial linkage, live FIR intake, "
             "multi-agent case work and proactive patrol.", bullet=True)
    para(tf, [("Trust by construction: ", WHITE, True),
              ("the model never invents numbers, the linkage ground truth is public, "
               "role scope is enforced in the query, and the audit log is "
               "tamper-evident. There is a red-team console in the product so anyone "
               "can test it.", BODY, False)], bullet=True)
    para(tf, [("Cross-district reach: ", WHITE, True),
              ("MO fingerprinting via multilingual embeddings finds serial offenders "
               "that siloed district records miss.", BODY, False)], bullet=True)
    para(tf, [("Built for real officers: ", WHITE, True),
              ("bilingual Kannada and English, voice or text, government-grade UI.",
               BODY, False)], bullet=True)
    para(tf, [("Where it makes a difference: ", WHITE, True),
              ("a constable at a rural station gets a statewide answer in Kannada in "
               "seconds; an SP spots a ring crossing Mandya into Bengaluru before the "
               "next strike; an IO gets a first-draft court-ready pack in minutes; a "
               "control room sees tomorrow's hotspot tonight.", BODY, False)],
         size=11, space=6)
    para(tf, [("USP:  ", ACCENT, True),
              ('"Ask a question, get a verified answer with evidence, and let the AI '
               'work the case end to end to a court-ready pack."', WHITE, False)],
         size=11.5, space=2)

    # Slide 3: List of features
    tf = add_body(S[3])
    feats = [
        ("Conversational floor", ": bilingual English and Kannada, voice or text, "
         "verified natural-language-to-SQL with an evidence drawer on every answer."),
        ("Serial Crime Linkage Engine", ": MO fingerprinting (multilingual embeddings "
         "plus structured features), weighted cosine, HDBSCAN, 180-day / 120-km "
         "filter. Discovers cross-district series cold."),
        ("Live FIR intake", ": file a new FIR in plain words, typed or dictated in "
         "Kannada. It is embedded at runtime and joins a matching series immediately."),
        ("AI Investigation Cell", ": six agents stream their reasoning and assemble a "
         "court-ready Investigation Pack, rendered to PDF by SmartBrowz."),
        ("CrimeGraph and Person 360", ": a knowledge graph over people and cases, plus "
         "one page holding a person's whole footprint, risk breakdown and associates."),
        ("Night Patrol and patrol plans", ": spike, series-growth and repeat-offender "
         "detectors produce ranked lead cards, then a deployment plan per district."),
        ("Trust Center", ": live metrics, a red-team console anyone can attack, and a "
         "hash-chained audit log that proves history has not been rewritten."),
        ("Governance", ": role scope enforced server-side (SHO, SP, SCRB, analyst with "
         "masked names), and no protected attributes in any model (ADR-9)."),
    ]
    for i, (h, rest) in enumerate(feats):
        para(tf, [(f"{i+1}. {h}", WHITE, True), (rest, BODY, False)], size=10.5,
             space=4, first=(i == 0))

    # Slide 4: Process flow / use-case
    tf = add_body(S[4])
    para(tf, "Officer (English or Kannada, voice or text) to React SPA to a FastAPI "
             "intent router:", size=12, color=WHITE, bold=True, first=True, space=8)
    para(tf, [("Question path: ", ACCENT, True),
              ("schema card and few-shots, then guardrails (SELECT-only, allowlisted "
               "tables, no file IO), the ADR-9 policy screen, role-scope injection, "
               "and self-repair. Answer plus evidence drawer.", BODY, False)],
         bullet=True)
    para(tf, [("Tool path: ", ACCENT, True),
              ("linkage, graph, forecast, hotspots, risk and similarity tools compute "
               "on DuckDB; the model narrates the verified result.", BODY, False)],
         bullet=True)
    para(tf, [("New FIR: ", ACCENT, True),
              ("the narrative is embedded at runtime by ONNX, written to Data Store "
               "and the mirror, linkage re-runs, and the officer is told which series "
               "it joined.", BODY, False)], bullet=True)
    para(tf, [("Investigation: ", ACCENT, True),
              ('"Investigate SH-07" starts the six-agent pipeline over SSE, ending in '
               "the Investigation Pack and its PDF.", BODY, False)], bullet=True)
    para(tf, [("Autonomous: ", ACCENT, True),
              ("Cron runs Night Patrol, which produces lead cards, a patrol plan and a "
               "spoken Kannada morning briefing.", BODY, False)], bullet=True)
    para(tf, [("Two data layers (ADR-1): ", WHITE, True),
              ("Catalyst Data Store is the system of record; an embedded DuckDB mirror "
               "serves sub-second analytics.", BODY, False)], size=11, space=2)

    # Slide 5: Wireframes / mock
    tf = add_body(S[5])
    para(tf, "Nine views in the live prototype (government-grade dark theme):",
         size=12, color=WHITE, bold=True, first=True, space=8)
    for h, rest in [
        ("Chat", ": bilingual answer, evidence drawer, voice input, scope badge."),
        ("Lead Feed", ": ranked Night-Patrol cards, patrol plan, spoken briefing."),
        ("Series", ": discovered series with codename, link explanations, replay and "
         "the counterfactual banner."),
        ("New FIR", ": file a case in plain words and watch it join a series."),
        ("Investigation Room", ": six agents streaming, then the assembled pack."),
        ("Person 360", ": one person's history, risk breakdown, network, associates."),
        ("CrimeGraph", ": interactive people and case network."),
        ("Trust Center", ": metrics, red-team console, audit-chain verification."),
        ("Audit", ": every action logged with its role."),
    ]:
        para(tf, [(h, WHITE, True), (rest, BODY, False)], bullet=True, size=11)
    para(tf, f"Live: {LIVE_URL}", size=10.5, color=ACCENT, space=2)

    # Slide 6: Architecture
    tf = add_body(S[6])
    flow = [
        "Officer (English / Kannada, voice and text)",
        "React SPA, served by the backend at /ui",
        "FastAPI on Catalyst AppSail, then the intent router",
        "  guardrails, ADR-9 policy, role scope, self-repair, and typed tools",
        "        (linkage, graph, forecast, hotspots, risk, similarity)",
        "DuckDB analytical mirror, with Catalyst Data Store as system of record",
        "ONNX MiniLM embeds new text at runtime, with no external model calls",
        "LLM adapter to Catalyst QuickML, serving GLM-4.7-Flash",
        "Six-agent Investigation Cell over SSE, then SmartBrowz for the pack PDF",
        "Cron drives Night Patrol, producing lead cards and the morning briefing",
    ]
    for i, line in enumerate(flow):
        indent = line.startswith(" ")
        para(tf, line.strip(), size=10.5, color=(MUTED if indent else BODY),
             first=(i == 0), space=3)
    para(tf, [("Deterministic tools compute, the model narrates (ADR-2). ", WHITE, True),
              ("No external inference API is ever called from the deployed app (ADR-4).",
               BODY, False)], size=10.5, space=2)

    # Slide 7: Technologies
    tf = add_body(S[7])
    for h, rest in [
        ("Backend", "Python 3.12, FastAPI, DuckDB, sqlglot, ONNX Runtime with "
                    "paraphrase-multilingual-MiniLM-L12-v2, scikit-learn and HDBSCAN, "
                    "NetworkX with Louvain, statsmodels SARIMA, sse-starlette."),
        ("Frontend", "React 18, Vite, Tailwind CSS, ECharts for charts and the network "
                     "graph, Leaflet with OpenStreetMap, Web Speech API for kn-IN voice."),
        ("LLM", "GLM-4.7-Flash served through Catalyst QuickML LLM Serving."),
        ("Packaging", "Docker custom AppSail runtime, with the scientific stack, the "
                      "embedding model and the SPA bundled in the image."),
        ("Quality", "86 tests, a 60-question bilingual eval harness, public linkage "
                    "ground truth, and a post-deploy verification gate."),
    ]:
        para(tf, [(f"{h}:  ", WHITE, True), (rest, BODY, False)], size=11.5, space=6,
             first=(h == "Backend"))

    # Slide 8: Catalyst services
    tf = add_body(S[8])
    para(tf, "Every part of ANVESHAK runs on Catalyst. The app reports its own "
             "platform status, so this is inspectable rather than asserted:",
         size=11.5, color=WHITE, bold=True, first=True, space=7)
    rows = [
        ("AppSail", "runs the FastAPI backend and serves the React app, on a custom "
                    "Docker runtime carrying the scientific stack and the embedding "
                    "model", GOOD, "LIVE"),
        ("QuickML (LLM Serving)", "serves GLM-4.7-Flash for question understanding "
                                  "and narration", GOOD, "LIVE"),
        ("API Gateway", "fronts the deployment", GOOD, "LIVE"),
        ("Data Store", "system of record for cases and new FIR intake; the "
                       "analytical mirror is rebuilt from it", ACCENT, "INTEGRATED"),
        ("SmartBrowz", "renders the court-ready Investigation Pack to PDF",
         ACCENT, "INTEGRATED"),
        ("Web Client Hosting", "hosts the built React bundle as a client component",
         MUTED, "CONFIGURED"),
        ("Authentication", "officer identity, feeding the role scopes the app "
                           "already enforces server-side", MUTED, "CONFIGURED"),
        ("Cron", "the overnight Night Patrol sweep, and cache warming",
         MUTED, "CONFIGURED"),
        ("NoSQL · Signals · Mail · Cache · Stratus",
         "graph snapshots, lead digests to district officers, shared state, and "
         "pack storage", MUTED, "CONFIGURED"),
    ]
    for name, purpose, colour, tag in rows:
        para(tf, [(f"{tag:11}", colour, True), (f"{name}: ", WHITE, True),
                  (purpose, BODY, False)], size=10, space=3)
    para(tf, [("No external inference API is ever called from the deployed app. ",
               WHITE, True),
              ("Open the Trust Center in the live app to see this table generated "
               "from the running process.", BODY, False)], size=10.5, space=2)

    # Slide 9: Estimated cost (optional)
    tf = add_body(S[9])
    para(tf, "The prototype runs within the Catalyst development tier.", size=12,
         color=WHITE, bold=True, first=True, space=8)
    para(tf, "AppSail: 1 GB memory, single instance, custom Docker runtime.", bullet=True)
    para(tf, "QuickML: GLM-4.7-Flash is a Zoho-hosted shared model, so there is no "
             "per-model deployment cost; it is billed per token at production scale.",
         bullet=True)
    para(tf, "Web Client Hosting, Data Store and Cron are serverless and usage-based.",
         bullet=True)
    para(tf, "Corpus embeddings are precomputed, and runtime embedding uses ONNX on "
             "CPU, so no GPU is required at any point.", bullet=True)

    # Slide 10: Snapshots
    _place_shots(S[10])

    # Slide 11: Performance / benchmarking
    tf = add_body(S[11])
    rows = [
        ("Dataset", "15,405 FIRs, 31 districts, 248 stations, 2023 to 2026 "
                    "(synthetic, seed 42)."),
        ("Natural language to SQL", "76.7% overall (English 82.2%, Kannada 60.0%) "
                                    "measured on the local 7B development model; "
                                    "gold-versus-gold 100%. A live harness re-measures "
                                    "against the deployed GLM."),
        ("Serial linkage (SH-07)", "precision 0.86, recall 12 of 14; the disjoint SP-2 "
                                   "series is correctly not merged."),
        ("Counterfactual", "the series becomes detectable at case 6; nine further "
                           "offences across three districts followed over 142 days."),
        ("Investigation Cell", "six agents complete in about 10 seconds, producing "
                               "8 ranked suspects."),
        ("Forecast (burglary)", "backtest MAE 1.83 against a seasonal-naive 1.75, "
                                "which is competitive and honestly reported."),
        ("Guardrails", "10 of 10 attack vectors blocked, re-tested live on every Trust "
                       "Center load."),
        ("Robustness", "thread-safe analytical layer: 60 concurrent requests, zero "
                       "errors; caches pre-warmed at startup."),
    ]
    for i, (h, rest) in enumerate(rows):
        para(tf, [(f"{h}:  ", WHITE, True), (rest, BODY, False)], size=10.5, space=4,
             first=(i == 0))
    para(tf, "Linkage ground truth is public (data_engine/planted/*.yaml), so these "
             "numbers are independently verifiable.", size=10, color=MUTED, space=2)

    # Slide 12: Links
    tf = add_body(S[12])
    para(tf, [("GitHub (public):  ", WHITE, True), (GITHUB, ACCENT, False)], size=13,
         first=True, space=10)
    para(tf, [("Deployed link (Catalyst):  ", WHITE, True), (LIVE_URL, ACCENT, False)],
         size=13, space=10)
    para(tf, [("Demo video (3 min):  ", WHITE, True),
              ("https://youtu.be/SpmZWo9keSI", ACCENT, False)],
         size=13, space=10)
    para(tf, f"API base: {API_URL}", size=10, color=MUTED, space=2)

    # Slide 13: Additional / future development
    tf = add_body(S[13])
    para(tf, "Roadmap to production:", size=12, color=WHITE, bold=True, first=True,
         space=8)
    for line in [
        "CCTNS integration: a live read connection plus identity resolution on name, "
        "parentage and date of birth, replacing the synthetic person key.",
        "Catalyst Auth issuing real officer identities into the scope layer that "
        "already enforces them server-side.",
        "Kannada tuning with hand-curated examples from real officer phrasing, which "
        "is our clearest measured gap.",
        "Human-in-the-loop learning: every analyst Confirm or Reject on a series "
        "becomes a label that improves precision.",
        "Photo intake: scan a paper FIR with a vision model, so the workflow starts "
        "where policing actually starts.",
        "A pilot with one district's live data.",
    ]:
        para(tf, line, bullet=True, size=11.5, color=BODY)

    removed = _delete_slide_by_text(prs, "Blank slide")
    prs.save(OUT)
    print(f"wrote {OUT}  ({len(prs.slides)} slides, blank removed={removed})")


if __name__ == "__main__":
    main()
